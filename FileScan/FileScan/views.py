from email.mime import audio
from http import client
from django.views.decorators.csrf import csrf_exempt
from os import listdir, remove
from django.http import HttpResponse, JsonResponse
import PyPDF2
import re
from pptx import Presentation
import moviepy.editor as mp
import random
import speech_recognition as speech_recog
import json
from google.cloud import speech_v1 as speech

config = dict(
    language_code="en-US",
    enable_automatic_punctuation=True,
    enable_word_time_offsets=True,
    sample_rate_hertz= 44100
)

@csrf_exempt
def getResources1(request):
   key = json.loads(request.body.decode('utf-8'))  
   keyWord = key['keyWord']
   ResSearch = ''
   pdfList = getLevel1PdfList(keyWord)
   presentationList = getLevel1PptxList(keyWord)
   videoKey = getVideoKeyWord(keyWord)
#    videoList = getLevel1VideoList(videoKey)
   listing = [  pdfList, presentationList ]
   return JsonResponse(listing, safe=False)

@csrf_exempt
def getResources2(request):
   key = json.loads(request.body.decode('utf-8'))  
   keyWord = key['keyWord']
   ResSearch = ''
   pdfList = getLevel2PdfList(keyWord)
   presentationList = getLevel2PptxList(keyWord)
   videoKey = getVideoKeyWord(keyWord)
#    videoList = getLevel2VideoList(videoKey)
   listing = [  pdfList, presentationList ]
   return JsonResponse(listing, safe=False)

@csrf_exempt
def getResources3(request):
   key = json.loads(request.body.decode('utf-8'))  
   keyWord = key['keyWord']
   ResSearch = ''
   pdfList = getLevel3PdfList(keyWord)
   presentationList = getLevel3PptxList(keyWord)
   videoKey = getVideoKeyWord(keyWord)
#    videoList = getLevel3VideoList(videoKey)
   listing = [  pdfList, presentationList ]
   return JsonResponse(listing, safe=False)

def getLevel1PdfList (word):
    pdfToStudy = []
    pdfList = listdir('Level1/pdf1/')

    for pdf in pdfList:
        object = PyPDF2.PdfFileReader('Level1/pdf1/'+pdf)
        NumPages = object.getNumPages()
        for i in range(0, NumPages):
            PageObj = object.getPage(i)
            Text = PageObj.extractText() 
            raw_search_string = r"\b" + str(word)
            print(raw_search_string)
            ResSearch = re.search(raw_search_string, Text, re.IGNORECASE)

            if(ResSearch is not None):
                pdfToStudy.insert(len(pdfToStudy)+1, pdf+" page no "+ str(i+1)+" Onwards ")
    return pdfToStudy

def getLevel1PptxList (word):
    slidesList = listdir('Level1/slides/')
    slidesToStudy = []

    for presentation in slidesList:
        prs = Presentation('Level1/slides/'+presentation)
        slideNo = 0
        for slide in prs.slides:
            slideNo = slideNo + 1
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    shape.text = shape.text.lower()
                    if str(word).casefold() in shape.text:
                        slidesToStudy.insert(len(slidesToStudy)+1, presentation+" slide no "+ str(slideNo)+" Onwards  ")
            # break
    return slidesToStudy

def getLevel1VideoList(word):
    videoToStudy = []
    videoList = listdir('Level1/video/')
    client = speech.SpeechClient.from_service_account_file('key.json')

    for video in videoList:
        
        my_clip = mp.VideoFileClip("Level1/video/"+video)
        my_clip.audio.write_audiofile("audio.mp3")
        with open("audio.mp3", 'rb') as f:
            mp3 = f.read()
        audio = speech.RecognitionAudio(content=mp3)
        aaa  = (client.recognize(
            config=config,
            audio= audio
        ))

        ssss = print_sentences(aaa, word)
        videoToStudy.append(str(video) + "  "+ str(ssss))

                
    return videoToStudy

def getLevel2PdfList (word):
    pdfToStudy = []
    pdfList = listdir('Level2/pdf/')

    for pdf in pdfList:
        object = PyPDF2.PdfFileReader('Level2/pdf/'+pdf)
        NumPages = object.getNumPages()
        for i in range(0, NumPages):
            PageObj = object.getPage(i)
            Text = PageObj.extractText() 
            raw_search_string = r"\b" + str(word)
            print(raw_search_string)
            ResSearch = re.search(raw_search_string, Text, re.IGNORECASE)

            if(ResSearch is not None):
                pdfToStudy.insert(len(pdfToStudy)+1, pdf+" page no "+ str(i+1)+" Onwards ")
    return pdfToStudy

def getLevel2PptxList (word):
    slidesList = listdir('Level2/slides/')
    slidesToStudy = []

    for presentation in slidesList:
        prs = Presentation('Level2/slides/'+presentation)
        slideNo = 0
        for slide in prs.slides:
            slideNo = slideNo + 1
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    shape.text = shape.text.lower()
                    if str(word).casefold() in shape.text:
                        slidesToStudy.insert(len(slidesToStudy)+1, presentation+" slide no "+ str(slideNo)+" Onwards  ")
            # break
    return slidesToStudy

def getLevel2VideoList(word):
    videoToStudy = []
    videoList = listdir('Level2/video/')
    client = speech.SpeechClient.from_service_account_file('key.json')

    for video in videoList:
        
        my_clip = mp.VideoFileClip("Level2/video/"+video)
        my_clip.audio.write_audiofile("audio.mp3")
        with open("audio.mp3", 'rb') as f:
            mp3 = f.read()
        audio = speech.RecognitionAudio(content=mp3)
        aaa  = (client.recognize(
            config=config,
            audio= audio
        ))

        ssss = print_sentences(aaa, word)
        videoToStudy.append(str(video) + "  "+ str(ssss))

                
    return videoToStudy

def getLevel3PdfList (word):
    pdfToStudy = []
    pdfList = listdir('Level3/pdf/')

    for pdf in pdfList:
        object = PyPDF2.PdfFileReader('Level3/pdf/'+pdf)
        NumPages = object.getNumPages()
        for i in range(0, NumPages):
            PageObj = object.getPage(i)
            Text = PageObj.extractText() 
            raw_search_string = r"\b" + str(word)
            print(raw_search_string)
            ResSearch = re.search(raw_search_string, Text, re.IGNORECASE)

            if(ResSearch is not None):
                pdfToStudy.insert(len(pdfToStudy)+1, pdf+" page no "+ str(i+1)+" Onwards ")
    return pdfToStudy

def getLevel3PptxList (word):
    slidesList = listdir('Level3/slides/')
    slidesToStudy = []

    for presentation in slidesList:
        prs = Presentation('Level3/slides/'+presentation)
        slideNo = 0
        for slide in prs.slides:
            slideNo = slideNo + 1
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    shape.text = shape.text.lower()
                    if str(word).casefold() in shape.text:
                        slidesToStudy.insert(len(slidesToStudy)+1, presentation+" slide no "+ str(slideNo)+" Onwards  ")
            # break
    return slidesToStudy

def getLevel3VideoList(word):
    videoToStudy = []
    videoList = listdir('Level3/video/')
    client = speech.SpeechClient.from_service_account_file('key.json')

    for video in videoList:
        
        my_clip = mp.VideoFileClip("Level3/video/"+video)
        my_clip.audio.write_audiofile("audio.mp3")
        with open("audio.mp3", 'rb') as f:
            mp3 = f.read()
        audio = speech.RecognitionAudio(content=mp3)
        aaa  = (client.recognize(
            config=config,
            audio= audio
        ))

        ssss = print_sentences(aaa, word)
        videoToStudy.append(str(video) + "  "+ str(ssss))

                
    return videoToStudy

def speech_to_text(config, audio):
    print("--------------------- hit------------------")
    client = speech.SpeechClient.from_service_account_file('key.json')
    response = client.recognize(config=config, audio=audio)
    print(response)
    # print_sentences(response)


def print_sentences(response, word):
    for result in response.results:
        best_alternative = result.alternatives[0]
        transcript = best_alternative.transcript
        confidence = best_alternative.confidence
        print("-" * 80)
        print(f"Transcript: {transcript}")
        print(f"Confidence: {confidence:.0%}")
        print("print_word_offsets   ", print_word_offsets(best_alternative, word))
        return print_word_offsets(best_alternative, word)


def print_word_offsets(alternative, words):
    for word in alternative.words:
        start_s = word.start_time.total_seconds()
        end_s = word.end_time.total_seconds()
        word = word.word
        print('word   ', word)
        if words == word :
            print("---------- hit -------------")
            return ( f"{start_s:>7.3f} | {end_s:>7.3f} | {word}")
        print(f"{start_s:>7.3f} | {end_s:>7.3f} | {word}")

def getVideoKeyWord (word):
    if word == 'OOP':
        videoKeyWord = 'inheritence'
    
    return videoKeyWord





